"""Agents: listing, chat, tuning (instructions + examples), document analysis."""

import io


async def _agents(auth):
    return (await auth.get("/api/v1/agents/")).json()["agents"]


async def _role_id(auth, role):
    return next(a["id"] for a in await _agents(auth) if a["role"] == role)


async def test_list_agents_seeded(auth):
    agents = await _agents(auth)
    assert len(agents) == 9
    roles = {a["role"] for a in agents}
    assert {"requirements", "ux", "architect", "developer", "tester", "devops"} <= roles


async def test_get_agent(auth):
    agents = await _agents(auth)
    aid = agents[0]["id"]
    r = await auth.get(f"/api/v1/agents/{aid}")
    assert r.status_code == 200
    assert r.json()["id"] == aid


async def test_agent_chat_persona_fallback(auth):
    aid = await _role_id(auth, "requirements")
    r = await auth.post(f"/api/v1/agents/{aid}/chat", json={"message": "How would you begin?"})
    assert r.status_code == 200
    body = r.json()
    assert body["source"] == "persona"  # no model configured in tests
    assert len(body["reply"]) > 0


async def test_instructions_get_set_reset(auth):
    aid = await _role_id(auth, "ux")
    g = await auth.get(f"/api/v1/agents/{aid}/instructions")
    assert g.status_code == 200
    assert g.json()["default_instructions"]
    assert g.json()["is_custom"] is False

    p = await auth.put(f"/api/v1/agents/{aid}/instructions", json={"instructions": "Be terse."})
    assert p.status_code == 200
    assert p.json()["instructions_custom"] is True
    assert p.json()["instructions"] == "Be terse."

    reset = await auth.put(f"/api/v1/agents/{aid}/instructions", json={"instructions": ""})
    assert reset.json()["instructions_custom"] is False


async def test_examples_save(auth):
    aid = await _role_id(auth, "architect")
    r = await auth.put(
        f"/api/v1/agents/{aid}/examples",
        json={"examples": [{"input": "scenario", "output": "The system shall do X."}]},
    )
    assert r.status_code == 200
    assert len(r.json()["examples"]) == 1
    # blank outputs are dropped
    r2 = await auth.put(f"/api/v1/agents/{aid}/examples", json={"examples": [{"output": "  "}]})
    assert r2.json()["examples"] == []


def _make_pptx() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Requirements Summary"
    tf = slide.placeholders[1].text_frame
    tf.text = "Users must be able to log in with email and password"
    for line in ["Export reports to PDF", "Send a notification when a task completes"]:
        tf.add_paragraph().text = line
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


async def test_analyze_pptx_to_requirements(auth):
    aid = await _role_id(auth, "requirements")
    files = {
        "file": (
            "summary.pptx",
            _make_pptx(),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    }
    r = await auth.post(f"/api/v1/agents/{aid}/analyze", files=files)
    assert r.status_code == 200
    body = r.json()
    assert body["source"] == "heuristic"  # no model in tests
    assert len(body["requirements"]) >= 2
    assert body["requirements"][0]["id"] == "FR-1"
    assert "system shall" in body["requirements"][0]["text"].lower()


async def test_analyze_rejects_unsupported_type(auth):
    aid = await _role_id(auth, "requirements")
    files = {"file": ("bad.exe", b"\x00\x01", "application/octet-stream")}
    r = await auth.post(f"/api/v1/agents/{aid}/analyze", files=files)
    assert r.status_code == 415


async def test_requirements_doc_download(auth):
    aid = await _role_id(auth, "requirements")
    r = await auth.post(
        f"/api/v1/agents/{aid}/requirements-doc",
        json={"title": "FRs", "requirements": ["The system shall allow login."]},
    )
    assert r.status_code == 200
    assert r.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )
    assert len(r.content) > 1000  # a real .docx binary
