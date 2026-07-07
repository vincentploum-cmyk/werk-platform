"""Extract plain text from uploaded requirements documents (pptx/pdf/docx/txt)."""

from __future__ import annotations

import io
import os


SUPPORTED = {".pptx", ".pdf", ".docx", ".txt", ".md"}


def supported(filename: str) -> bool:
    return os.path.splitext(filename.lower())[1] in SUPPORTED


def extract_text(filename: str, content: bytes) -> str:
    """Return plain text extracted from the file, dispatched by extension."""
    ext = os.path.splitext(filename.lower())[1]
    if ext == ".pptx":
        return _from_pptx(content)
    if ext == ".pdf":
        return _from_pdf(content)
    if ext == ".docx":
        return _from_docx(content)
    if ext in (".txt", ".md"):
        return content.decode("utf-8", errors="replace")
    raise ValueError(f"Unsupported file type: {ext or '(none)'}")


def _from_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        slide_lines.append(" | ".join(cells))
        if slide_lines:
            lines.append(f"[Slide {i}]")
            lines.extend(slide_lines)
    return "\n".join(lines)


def _from_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages = []
    for page in reader.pages:
        try:
            pages.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n".join(pages)


def _from_docx(content: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(content))
    lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)
